# Import libraries
import streamlit as st
from langchain.llms import OpenAI
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import Chroma
from langchain.chains import RetrievalQA
from PyPDF2 import PdfReader
from io import BytesIO
from langchain.text_splitter import RecursiveCharacterTextSplitter
from pptx import Presentation
from docx import Document
from langchain.text_splitter import TokenTextSplitter

# Install sqlite3 module
__import__('pysqlite3')
import sys
sys.modules['sqlite3'] = sys.modules.pop('pysqlite3')

def read_docx(file_path):
    document = Document(file_path)
    text = ''
    for paragraph in document.paragraphs:
        text += paragraph.text
    return text

def read_pptx(file_path):
    presentation = Presentation(file_path)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text
    return text

def generate_response(uploaded_file, openai_api_key, query_text):
    # Load document if file is uploaded
    if uploaded_file is not None:
        # Get the file extension
        file_type = uploaded_file.type

        # Use PyPDF2 to read the PDF file
        if file_type == "application/pdf":
            pdf = PdfReader(BytesIO(uploaded_file.getvalue()))
            # Extract text from each page and add it to the text string
            text = ''
            for page in pdf.pages:
                text += page.extract_text()
                
        # Read the TXT file
        elif file_type == "text/plain":
            text = uploaded_file.getvalue().decode()
            
        elif "presentationml" in file_type: # pptx file
            text = read_pptx(uploaded_file)
            
        elif "wordprocessingml" in file_type: # docx file
            text = read_docx(uploaded_file)

        documents = [text]
        # Split documents into chunks
        text_splitter = TokenTextSplitter(
            chunk_size=800, 
            chunk_overlap=10
            )
        # text_splitter = RecursiveCharacterTextSplitter(
        #     # In English, one token is roughly equivalent to 4 characters1. 
        #     # So, approximately 3500 tokens would be around 14,000 characters.
        #     chunk_size=3500, 
        #     chunk_overlap=20,
        #     length_function=len,
        #     is_separator_regex=False,
        # )
        texts = text_splitter.create_documents(documents)

        # Select embeddings
        embeddings = OpenAIEmbeddings(openai_api_key=openai_api_key)
        # Create a vectorstore from documents
        # Chroma vectorstoreï¼šdesigned for short documents, such as tweets and product descriptions. It is not ideal for long documents, such as PDF files and presentations.
        # HNSW vectorstore: This vectorstore is designed for long documents, and it is more efficient than the Chroma vectorstore.
        db = Chroma.from_documents(texts, embeddings)
        
        # Create retriever interface
        retriever = db.as_retriever()
        
        # Create QA chain
        qa = RetrievalQA.from_chain_type(llm=OpenAI(openai_api_key=openai_api_key), chain_type='stuff', retriever=retriever)
        return qa.run(query_text)

# Page title
st.set_page_config(page_title='ðŸ¦œðŸ”— Ask the Doc App')
st.title('ðŸ¦œðŸ”— Ask the Doc App')

# File upload
uploaded_file = st.file_uploader('Upload a PDF, PPTX, DOCX or TXT file', type=["pdf", "txt", "pptx", "docx"])
# Query text
query_text = st.text_input('Enter your question:', placeholder = 'Please provide a short summary.', disabled=not uploaded_file)

openai_api_key = "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"
# Form input and query
result = []
with st.form('myform', clear_on_submit=True):
    # openai_api_key = st.text_input('OpenAI API Key', type='password', disabled=not (uploaded_file and query_text))
    submitted = st.form_submit_button('Submit', disabled=not(uploaded_file and query_text))
    if submitted and openai_api_key.startswith('sk-'):
        with st.spinner('Calculating...'):
            response = generate_response(uploaded_file, openai_api_key, query_text)
            result.append(response)
            # del openai_api_key

if len(result):
    st.info(response)
